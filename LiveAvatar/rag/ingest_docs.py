import os
import json
import uuid
import re
import hashlib
import time
from pathlib import Path
from typing import Iterator, List, Tuple

import fitz  # PyMuPDF
from pptx import Presentation
from dotenv import load_dotenv
from tqdm import tqdm

from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient

from embedding_client import get_embeddings

# ----------------------------
# Config
# ----------------------------
load_dotenv()

SEARCH_SERVICE = os.getenv("AZURE_SEARCH_SERVICE")
SEARCH_KEY = os.getenv("AZURE_SEARCH_ADMIN_KEY")
INDEX_NAME = os.getenv("AZURE_SEARCH_INDEX", "training-index")
DOCS_ROOT = Path(os.getenv("DOCS_ROOT", "./docs")).resolve()

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}

MAX_FILE_MB = int(os.getenv("MAX_FILE_MB", "200"))
MAX_PDF_PAGES = int(os.getenv("MAX_PDF_PAGES", "300"))

CHUNK_MAX_CHARS = int(os.getenv("CHUNK_MAX_CHARS", "2000"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "150"))

EMBED_BATCH_SIZE = int(os.getenv("EMBED_BATCH_SIZE", "32"))
UPLOAD_BATCH_SIZE = int(os.getenv("UPLOAD_BATCH_SIZE", "200"))
SLEEP_BETWEEN_BATCHES_SEC = float(os.getenv("SLEEP_BETWEEN_BATCHES_SEC", "0.0"))

MAX_DOC_CHARS = int(os.getenv("MAX_DOC_CHARS", "300000"))

CHECKPOINT_PATH = Path(os.getenv("CHECKPOINT_PATH", "./ingest_checkpoint.json")).resolve()

if not SEARCH_SERVICE or not SEARCH_KEY:
    raise RuntimeError("Missing AZURE_SEARCH_SERVICE or AZURE_SEARCH_ADMIN_KEY")
if not DOCS_ROOT.exists():
    raise RuntimeError(f"DOCS_ROOT does not exist: {DOCS_ROOT}")

search_client = SearchClient(
    endpoint=SEARCH_SERVICE.rstrip("/"),
    index_name=INDEX_NAME,
    credential=AzureKeyCredential(SEARCH_KEY),
)

# ----------------------------
# Filters / cleanup
# ----------------------------
SKIP_PATH_SUBSTRINGS = [
    "Training Certificates",
    "Certificate Templates",
    "/Forms/",
]

BOILERPLATE_PATTERNS = [
    r"WWW\.GSGLI\.COM",
    r"\bTF(F|B)\s*#\s*\d+",
    r"state of the art materials",
]

def should_skip_path(p: Path) -> bool:
    s = str(p)
    return any(x in s for x in SKIP_PATH_SUBSTRINGS)

def clean_text(text: str) -> str:
    lines = text.splitlines()
    cleaned = []
    for ln in lines:
        s = ln.strip()
        if not s:
            continue
        if any(re.search(pat, s, re.IGNORECASE) for pat in BOILERPLATE_PATTERNS):
            continue
        cleaned.append(s)
    return "\n".join(cleaned)

def content_hash(s: str) -> str:
    return hashlib.sha1(s.encode("utf-8", errors="ignore")).hexdigest()

# ----------------------------
# Checkpoint
# ----------------------------
def load_checkpoint() -> set:
    if CHECKPOINT_PATH.exists():
        try:
            data = json.loads(CHECKPOINT_PATH.read_text(encoding="utf-8"))
            return set(data.get("completed_files", []))
        except Exception:
            return set()
    return set()

def save_checkpoint(completed_files: set):
    CHECKPOINT_PATH.write_text(
        json.dumps({"completed_files": sorted(list(completed_files))}, indent=2),
        encoding="utf-8",
    )

# ----------------------------
# Metadata
# ----------------------------
def infer_metadata(root: Path, file_path: Path) -> Tuple[str, str, str]:
    rel = file_path.relative_to(root)
    parts = rel.parts
    # DOCS_ROOT / <course> / <module> / ...
    course_id = parts[0] if len(parts) >= 2 else "general"
    module_id = parts[1] if len(parts) >= 3 else "misc"
    doc_id = file_path.stem
    return doc_id, course_id, module_id

# ----------------------------
# Chunking (simple + safe)
# ----------------------------
def iter_chunks(text: str, max_chars: int, overlap: int) -> Iterator[str]:
    # normalize whitespace
    text = " ".join(text.split())
    n = len(text)
    if n == 0:
        return
    start = 0
    while start < n:
        end = min(start + max_chars, n)
        yield text[start:end]
        if end == n:
            break
        start = max(0, end - overlap)

# ----------------------------
# Extractors
# ----------------------------
def extract_text_pdf_pages(path: Path, max_pages: int) -> Iterator[tuple[int, str]]:
    doc = fitz.open(path)
    try:
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            yield i, page.get_text()
    finally:
        doc.close()

import zipfile
from xml.etree import ElementTree as ET

def extract_text_docx_xml(path: Path) -> str:
    # Reads word/document.xml and concatenates all w:t nodes.
    with zipfile.ZipFile(path) as z:
        xml = z.read("word/document.xml")
    root = ET.fromstring(xml)
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

    texts = []
    for t in root.findall(".//w:t", ns):
        if t.text:
            texts.append(t.text)

    out = " ".join(texts)
    out = re.sub(r"\s+", " ", out).strip()
    return out

def extract_text_pptx(path: Path) -> str:
    prs = Presentation(path)
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                out.append(shape.text)
    return "\n".join(out)

# ----------------------------
# Upload
# ----------------------------
def upload_docs(docs: List[dict]):
    for i in range(0, len(docs), UPLOAD_BATCH_SIZE):
        search_client.upload_documents(docs[i : i + UPLOAD_BATCH_SIZE])

# ----------------------------
# Ingest single file
# ----------------------------
def ingest_file(file_path: Path) -> int:
    size_mb = file_path.stat().st_size / (1024 * 1024)
    if size_mb > MAX_FILE_MB:
        print(f"⚠️ Skipping (too large {size_mb:.1f}MB > {MAX_FILE_MB}MB): {file_path}")
        return 0

    doc_id, course_id, module_id = infer_metadata(DOCS_ROOT, file_path)
    ext = file_path.suffix.lower()
    source_type = ext.lstrip(".")
    seen_hashes = set()

    total_uploaded = 0
    chunk_num_global = 0
    pending: List[tuple[int, int, str]] = []  # (page_num, chunk_num, text)

    def flush_pending(chunks: List[tuple[int, int, str]]) -> int:
        if not chunks:
            return 0
        texts = [t for (_, _, t) in chunks]
        vectors = get_embeddings(texts)

        docs = []
        for (page_num, chunk_num, t), vec in zip(chunks, vectors):
            h = content_hash(t)
            docs.append(
                {
                    "id": str(uuid.uuid4()),
                    "content": t,
                    "content_vector": vec,
                    "doc_id": doc_id,
                    "course_id": course_id,
                    "module_id": module_id,
                    "path": str(file_path),
                    "source_type": source_type,
                    "page_num": int(page_num),
                    "chunk_num": int(chunk_num),
                    "content_hash": h,
                }
            )

        upload_docs(docs)
        if SLEEP_BETWEEN_BATCHES_SEC > 0:
            time.sleep(SLEEP_BETWEEN_BATCHES_SEC)
        return len(docs)

    # ---- PDF ----
    if ext == ".pdf":
        for page_num, page_text in extract_text_pdf_pages(file_path, MAX_PDF_PAGES):
            page_text = clean_text(page_text)
            page_text = " ".join(page_text.split())
            if not page_text:
                continue

            # optional cap per-page to avoid weird PDF extraction issues
            if len(page_text) > MAX_DOC_CHARS:
                page_text = page_text[:MAX_DOC_CHARS]

            for chunk in iter_chunks(page_text, CHUNK_MAX_CHARS, CHUNK_OVERLAP):
                chunk = chunk.strip()
                if not chunk:
                    continue
                h = content_hash(chunk)
                if h in seen_hashes:
                    continue
                seen_hashes.add(h)

                pending.append((page_num, chunk_num_global, chunk))
                chunk_num_global += 1

                if len(pending) >= EMBED_BATCH_SIZE:
                    total_uploaded += flush_pending(pending)
                    pending = []

        total_uploaded += flush_pending(pending)
        return total_uploaded

    # ---- DOCX / PPTX ----
    print("   - extracting text…")
    if ext == ".docx":
        text = extract_text_docx_xml(file_path)
    elif ext == ".pptx":
        text = extract_text_pptx(file_path)
    else:
        return 0

    text = clean_text(text)
    text = " ".join(text.split())

    if not text:
        print("   - no usable text after cleaning")
        return 0

    if len(text) > MAX_DOC_CHARS:
        print(f"   - huge text ({len(text)} chars). Truncating to {MAX_DOC_CHARS}.")
        text = text[:MAX_DOC_CHARS]

    for chunk in iter_chunks(text, CHUNK_MAX_CHARS, CHUNK_OVERLAP):
        chunk = chunk.strip()
        if not chunk:
            continue
        h = content_hash(chunk)
        if h in seen_hashes:
            continue
        seen_hashes.add(h)

        pending.append((-1, chunk_num_global, chunk))
        chunk_num_global += 1

        if len(pending) >= EMBED_BATCH_SIZE:
            total_uploaded += flush_pending(pending)
            pending = []

    total_uploaded += flush_pending(pending)
    return total_uploaded

# ----------------------------
# File discovery
# ----------------------------
def iter_files(root: Path) -> List[Path]:
    out = []
    for p in root.rglob("*"):
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS:
            if should_skip_path(p):
                continue
            out.append(p)
    # smallest first so you see progress quickly
    out.sort(key=lambda p: p.stat().st_size)
    return out

# ----------------------------
# Main
# ----------------------------
def main():
    completed = load_checkpoint()

    # Write checkpoint immediately so you can confirm path/permissions.
    save_checkpoint(completed)
    print("✅ checkpoint write ok")

    print(f"✅ DOCS_ROOT: {DOCS_ROOT}")

    files = iter_files(DOCS_ROOT)

    print(f"✅ Files to ingest: {len(files)}")
    print(f"✅ Index: {INDEX_NAME}")
    print(f"✅ Checkpoint: {CHECKPOINT_PATH}\n")

    total_chunks = 0

    for f in tqdm(files, desc="Ingesting files"):
        key = str(f)
        if key in completed:
            continue

        print(f"\n➡️ Processing: {f} ({f.stat().st_size/1024:.0f} KB)")
        try:
            uploaded = ingest_file(f)
            total_chunks += uploaded
            completed.add(key)
            save_checkpoint(completed)
            print(f"✅ Uploaded {uploaded} chunks (running total {total_chunks})")
        except Exception as e:
            print(f"⚠️ Error ingesting {f}: {e}")
            print("   Not marking completed — fix and rerun.")
            continue

    print(f"\n✅ Uploaded total chunks: {total_chunks}")
    print(f"✅ Completed files: {len(completed)}")

if __name__ == "__main__":
    main()
